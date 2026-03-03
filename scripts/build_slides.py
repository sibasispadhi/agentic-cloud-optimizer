#!/usr/bin/env python3
"""Build PPTX + PDF backups from the offline Reveal.js deck.

Not trying to perfectly clone Reveal layout (YAGNI). This is a *bulletproof fallback*
format for conference A/V.

Inputs:
- presentation/devnexus-2026-slides-OFFLINE.html

Outputs:
- presentation/devnexus-2026-slides.pptx
- presentation/devnexus-2026-slides.pdf
"""

from __future__ import annotations

import base64
import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas

import qrcode


WM_BLUE = RGBColor(0x00, 0x53, 0xE2)  # #0053e2
WM_SPARK = RGBColor(0xFF, 0xC2, 0x20)  # #ffc220
WM_GRAY_160 = RGBColor(0x2E, 0x2F, 0x32)
WM_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: list[str]


def _clean_text(txt: str) -> str:
    t = re.sub(r"\s+", " ", txt or "").strip()
    return t


def parse_reveal_slides(html_path: Path) -> list[SlideSpec]:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "lxml")

    specs: list[SlideSpec] = []
    for sec in soup.find_all("section"):
        # Title: prefer h1 then h2, else skip
        h = sec.find(["h1", "h2"])
        if not h:
            continue
        title = _clean_text(h.get_text(" "))

        bullets: list[str] = []
        for ul in sec.find_all("ul"):
            for li in ul.find_all("li"):
                bullets.append(_clean_text(li.get_text(" ")))

        # keep it short; fallback deck should be readable
        bullets = [b for b in bullets if b][:8]
        specs.append(SlideSpec(title=title, bullets=bullets))

    # Reveal decks sometimes include nested sections; ensure we keep 19-ish slides.
    # If it’s more, we’ll keep first 19. If less, we still proceed.
    return specs[:19]


def make_qr_png_bytes(url: str) -> bytes:
    qr = qrcode.QRCode(version=2, box_size=8, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#0053e2", back_color="white")
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def build_pptx(specs: list[SlideSpec], out_path: Path, repo_url: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]  # title slide
    content_layout = prs.slide_layouts[1]  # title + content

    for idx, s in enumerate(specs, start=1):
        layout = title_layout if idx == 1 else content_layout
        slide = prs.slides.add_slide(layout)

        # Background: white
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WM_WHITE

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = s.title
            tf = slide.shapes.title.text_frame
            p = tf.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(38 if idx == 1 else 34)
            p.font.color.rgb = WM_BLUE

        # Subtitle / bullets
        if idx == 1:
            # Add a small subtitle block + repo url
            left, top, width, height = Inches(1.0), Inches(4.0), Inches(11.3), Inches(2.5)
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            tf.clear()

            p1 = tf.paragraphs[0]
            p1.text = "DevNexus 2026 · Offline-safe backup deck"
            p1.font.name = "Segoe UI"
            p1.font.size = Pt(20)
            p1.font.color.rgb = WM_GRAY_160

            p2 = tf.add_paragraph()
            p2.text = repo_url
            p2.font.name = "Consolas"
            p2.font.size = Pt(16)
            p2.font.color.rgb = WM_GRAY_160

            # Accent bar
            bar = slide.shapes.add_shape(
                1, Inches(0), Inches(7.25), prs.slide_width, Inches(0.25)  # MSO_AUTO_SHAPE_TYPE=1 rect
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = WM_SPARK
            bar.line.color.rgb = WM_SPARK

        else:
            body = slide.shapes.placeholders[1].text_frame if len(slide.shapes.placeholders) > 1 else None
            if body is not None:
                body.clear()
                if not s.bullets:
                    p = body.paragraphs[0]
                    p.text = "(See HTML deck for visuals)"
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(22)
                    p.font.color.rgb = WM_GRAY_160
                else:
                    for i, b in enumerate(s.bullets):
                        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                        p.text = b
                        p.level = 0
                        p.font.name = "Segoe UI"
                        p.font.size = Pt(22)
                        p.font.color.rgb = WM_GRAY_160

        # slide number footer (except title)
        if idx != 1:
            footer = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(1.0), Inches(0.3))
            tf = footer.text_frame
            tf.text = f"{idx}/{len(specs)}"
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].font.name = "Segoe UI"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = WM_BLUE

    # Ensure QR code exists on final slide (append if needed)
    qr_png = make_qr_png_bytes(repo_url)
    slide = prs.slides.add_slide(content_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WM_BLUE

    slide.shapes.title.text = "Thank You!"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = WM_WHITE
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # QR image
    img_stream = BytesIO(qr_png)
    slide.shapes.add_picture(img_stream, Inches(1.0), Inches(2.3), width=Inches(2.2), height=Inches(2.2))

    # URL text
    tb = slide.shapes.add_textbox(Inches(3.5), Inches(2.4), Inches(9.5), Inches(2.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Scan for source code"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WM_SPARK

    p2 = tf.add_paragraph()
    p2.text = "github.com/sibasispadhi/agentic-cloud-optimizer"
    p2.font.name = "Consolas"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WM_WHITE

    prs.save(str(out_path))


def build_pdf(specs: list[SlideSpec], out_path: Path, repo_url: str) -> None:
    # Simple landscape letter; projector-safe.
    w, h = landscape((11 * inch, 8.5 * inch))
    c = canvas.Canvas(str(out_path), pagesize=(w, h))

    def draw_header(title: str) -> None:
        c.setFillColorRGB(0, 0.325, 0.886)  # WM blue
        c.setFont("Helvetica-Bold", 26)
        c.drawString(0.7 * inch, h - 1.0 * inch, title)
        c.setStrokeColorRGB(1, 0.761, 0.125)  # spark
        c.setLineWidth(6)
        c.line(0.7 * inch, h - 1.1 * inch, w - 0.7 * inch, h - 1.1 * inch)

    for i, s in enumerate(specs, start=1):
        draw_header(s.title)
        c.setFillColorRGB(0.18, 0.184, 0.196)
        c.setFont("Helvetica", 16)

        y = h - 1.6 * inch
        for b in s.bullets[:10]:
            c.drawString(0.9 * inch, y, f"• {b}")
            y -= 0.35 * inch

        c.setFillColorRGB(0, 0.325, 0.886)
        c.setFont("Helvetica", 10)
        c.drawRightString(w - 0.7 * inch, 0.4 * inch, f"{i}/{len(specs)}")
        c.showPage()

    # final QR slide
    draw_header("Thank You!")
    qr_png = make_qr_png_bytes(repo_url)
    img_path = out_path.parent / "_tmp_qr.png"
    img_path.write_bytes(qr_png)
    c.drawImage(str(img_path), 0.9 * inch, h - 4.6 * inch, width=2.2 * inch, height=2.2 * inch, mask='auto')
    img_path.unlink(missing_ok=True)

    c.setFillColorRGB(1, 0.761, 0.125)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(3.4 * inch, h - 2.7 * inch, "Scan for source code")
    c.setFillColorRGB(0.18, 0.184, 0.196)
    c.setFont("Courier", 14)
    c.drawString(3.4 * inch, h - 3.2 * inch, repo_url)
    c.showPage()

    c.save()


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    html_path = repo_root / "presentation" / "devnexus-2026-slides-OFFLINE.html"
    if not html_path.exists():
        raise SystemExit(f"Missing input deck: {html_path}")

    specs = parse_reveal_slides(html_path)
    if not specs:
        raise SystemExit("No slides found in HTML. Is the deck valid?")

    out_pptx = repo_root / "presentation" / "devnexus-2026-slides.pptx"
    out_pdf = repo_root / "presentation" / "devnexus-2026-slides.pdf"
    repo_url = "https://github.com/sibasispadhi/agentic-cloud-optimizer"

    build_pptx(specs, out_pptx, repo_url)
    build_pdf(specs, out_pdf, repo_url)

    print(f"✅ wrote {out_pptx}")
    print(f"✅ wrote {out_pdf}")


if __name__ == "__main__":
    main()
